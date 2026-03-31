#!/usr/bin/env python3
"""
Generate Gauss-Markov PPTX slides following Professor Jiang's style.
Style: Gold (#FFC000) emphasis + Deep Blue (#0563C1) titles + Gray-Blue (#44546A) body
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import matplotlib.pyplot as plt
import matplotlib
matplotlib.use('Agg')
import numpy as np
import os
import tempfile

# Color scheme
GOLD = RGBColor(0xFF, 0xC0, 0x00)      # #FFC000 emphasis
DEEP_BLUE = RGBColor(0x05, 0x63, 0xC1) # #0563C1 titles
GRAY_BLUE = RGBColor(0x44, 0x54, 0x6A) # #44546A body text
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xA5, 0xA5, 0xA5)

# Font sizes
TITLE_SIZE = Pt(40)
SUBTITLE_SIZE = Pt(28)
BODY_SIZE = Pt(24)
SMALL_SIZE = Pt(18)

# Output path
OUTPUT_PATH = "/Users/yueyh/Projects/aca-workflow/notes/applied-linear-regression/gauss-markov-slides.pptx"

def create_formula_image(formula, filename):
    """Render formula to image using matplotlib."""
    fig, ax = plt.subplots(figsize=(6, 0.8))
    ax.text(0.5, 0.5, f'${formula}$', fontsize=16, ha='center', va='center')
    ax.axis('off')
    fig.savefig(filename, bbox_inches='tight', transparent=True, dpi=150)
    plt.close(fig)
    return filename

def add_geometric_decoration(slide, prs):
    """Add geometric rectangle decoration to slide (88% style)."""
    # Add a subtle rectangle decoration
    left = Inches(0)
    top = Inches(0)
    width = Inches(0.15)
    height = prs.slide_height
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()

def set_title_style(title_shape):
    """Apply title formatting."""
    title_shape.font.size = TITLE_SIZE
    title_shape.font.bold = True
    title_shape.font.color.rgb = DEEP_BLUE

def add_bullet_text(shape, bullets, font_size=BODY_SIZE):
    """Add bullet points to a text frame."""
    tf = shape.text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = font_size
        p.font.color.rgb = GRAY_BLUE
        p.level = 0

def add_slide_header(slide, label, label_color=GOLD):
    """Add a small colored label in top-left corner."""
    left = Inches(0.3)
    top = Inches(0.2)
    width = Inches(1.5)
    height = Inches(0.35)

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = label_color
    shape.line.fill.background()

    tf = shape.text_frame
    tf.paragraphs[0].text = label
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

def create_slide_content(slide, title, bullets, prs, label=None, formula_images=None):
    """Generic content slide creation."""
    add_geometric_decoration(slide, prs)

    # Title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_shape.text_frame.paragraphs[0].text = title
    set_title_style(title_shape.text_frame.paragraphs[0])

    # Label if provided
    if label:
        add_slide_header(slide, label)

    # Content area - bullet points
    content_top = Inches(1.5) if label else Inches(1.3)
    content_height = Inches(4.5) if not formula_images else Inches(3.5)
    content_shape = slide.shapes.add_textbox(Inches(0.5), content_top, Inches(9), content_height)
    add_bullet_text(content_shape, bullets)

    # Add formula images if provided
    if formula_images:
        img_top = Inches(5.0)
        for i, img_path in enumerate(formula_images):
            if os.path.exists(img_path):
                slide.shapes.add_picture(img_path, Inches(1.0 + i*4), img_top, width=Inches(4))

def create_cover_slide(prs):
    """Slide 1: Cover page."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Add decorative rectangle on left
    left = Inches(0)
    top = Inches(0)
    width = Inches(0.3)
    height = prs.slide_height
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()

    # Title
    title_shape = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(8.5), Inches(1.2))
    title_tf = title_shape.text_frame
    title_tf.paragraphs[0].text = "Gauss-Markov Theorem"
    title_tf.paragraphs[0].font.size = Pt(48)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = DEEP_BLUE

    # Subtitle
    subtitle_shape = slide.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(8), Inches(0.6))
    subtitle_tf = subtitle_shape.text_frame
    subtitle_tf.paragraphs[0].text = "应用回归分析"
    subtitle_tf.paragraphs[0].font.size = Pt(28)
    subtitle_tf.paragraphs[0].font.color.rgb = GRAY_BLUE

    # Course info
    info_shape = slide.shapes.add_textbox(Inches(0.8), Inches(4.6), Inches(8), Inches(0.5))
    info_tf = info_shape.text_frame
    info_tf.paragraphs[0].text = "Chapter 04 | Peng Ding, A First Course in Causal Inference"
    info_tf.paragraphs[0].font.size = Pt(18)
    info_tf.paragraphs[0].font.color.rgb = LIGHT_GRAY

def create_motivation_slide_1(prs):
    """Slide 2: Motivation - Why Gauss-Markov?"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bullets = [
        "OLS 估计凭什么\"最好\"？",
        "有没有比 OLS 更好的估计？",
        "第 3 章 OLS 是纯代数运算",
        "没有随机假设谈不上\"最优\""
    ]
    create_slide_content(slide, "为什么需要 Gauss-Markov？", bullets, prs, "Motivation")
    # Add question box
    q_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6), Inches(4.5), Inches(3.5), Inches(1.2))
    q_shape.fill.solid()
    q_shape.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
    q_shape.line.fill.background()
    tf = q_shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "\"OLS 凭什么最好？\""
    p.font.size = Pt(18)
    p.font.italic = True
    p.font.color.rgb = DEEP_BLUE

def create_motivation_slide_2(prs):
    """Slide 3: From Algebra to Statistics."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bullets = [
        "第 3 章 OLS 是纯代数运算",
        "没有随机假设谈不上\"最优\"",
        "Gauss-Markov 模型给出了讨论统计性质的基础",
        "需要误差的前二阶矩：均值、方差、协方差"
    ]
    create_slide_content(slide, "从代数到统计：引入随机假设", bullets, prs, "Motivation")

def create_transition_slide(prs):
    """Slide 4: Chapter Roadmap."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bullets = [
        "1. Gauss-Markov 模型假设",
        "2. OLS 估计量的性质（均值、方差）",
        "3. Gauss-Markov 定理（核心）",
        "4. 定理的直观解释与证明思路"
    ]
    create_slide_content(slide, "章节路线图", bullets, prs, "Roadmap")

def create_assumptions_slide(prs):
    """Slide 5: Gauss-Markov Assumptions."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bullets = [
        "线性形式: $Y = X\\beta + \\varepsilon$",
        "设计矩阵: $X$ 固定且列线性无关",
        "误差均值: $E(\\varepsilon) = 0$",
        "误差协方差: $\\cov(\\varepsilon) = \\sigma^2 I_n$",
        "（同方差、不相关）",
        "未知参数为 $(\\beta, \\sigma^2)$"
    ]
    create_slide_content(slide, "Gauss-Markov 模型假设", bullets, prs, "Assumption 4.1")

def create_individual_level_slide(prs):
    """Slide 6: Individual level perspective."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bullets = [
        "每个观测: $y_i = x_i^\\top \\beta + \\varepsilon_i$",
        "误差均值 0: $E(\\varepsilon_i) = 0$",
        "误差方差 $\\sigma^2$（同方差）",
        "不相关: $\\cov(\\varepsilon_i, \\varepsilon_j) = 0, i \\neq j$"
    ]
    create_slide_content(slide, "个体水平视角", bullets, prs, "概念")

def create_homoskedasticity_slide(prs):
    """Slide 7: Homoskedasticity importance."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bullets = [
        "homoskedasticity = 相同方差",
        "词源学：k 更好地表示 variance 含义",
        "(McCulloch 1985)",
        "异方差情形需要加权最小二乘",
        "（见第 19 章）"
    ]
    create_slide_content(slide, "同方差假设的含义与重要性", bullets, prs, "概念")

def create_ols_form_slide(prs):
    """Slide 8: OLS formula."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bullets = [
        "$\\hat{\\beta} = (X^\\top X)^{-1} X^\\top Y$",
        "$\\hat{\\beta}$ 是 $Y$ 的线性函数",
        "仅依赖矩阵运算",
        "令 $A = (X^\\top X)^{-1} X^\\top$，则 $\\hat{\\beta} = AY$",
        "$A$ 不依赖 $Y$，所以 OLS 是线性估计量"
    ]
    create_slide_content(slide, "OLS 估计量：矩阵形式", bullets, prs, "Theorem 4.1 前")

def create_ols_mean_var_slide(prs):
    """Slide 9: OLS mean and variance theorem."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bullets = [
        "Theorem 4.1 (无偏性): $E(\\hat{\\beta}) = \\beta$",
        "协方差矩阵: $\\cov(\\hat{\\beta}) = \\sigma^2 (X^\\top X)^{-1}$",
        "证明：利用 $E(Y) = X\\beta$ 和 $\\cov(Y) = \\sigma^2 I_n$"
    ]
    create_slide_content(slide, "OLS 估计量的均值与方差", bullets, prs, "Theorem 4.1")

def create_gauss_markov_theorem_slide(prs):
    """Slide 10: Gauss-Markov Theorem (core)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bullets = [
        "Gauss-Markov Theorem 4.2: $\\hat{\\beta}$ 是 BLUE",
        "Best Linear Unbiased Estimator",
        "条件 C1: $\\tilde{\\beta} = AY$ 对 $Y$ 线性，$A$ 不依赖 $Y$",
        "条件 C2: $E(\\tilde{\\beta}) = \\beta$（对所有 $\\beta$ 无偏）",
        "核心不等式: $\\cov(\\tilde{\\beta}) \\succeq \\cov(\\hat{\\beta})$"
    ]
    create_slide_content(slide, "Gauss-Markov 定理（核心）", bullets, prs, "Theorem 4.2")

def create_blue_meaning_slide(prs):
    """Slide 11: BLUE meaning."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bullets = [
        "$\\cov(\\tilde{\\beta}) \\succeq \\cov(\\hat{\\beta})$",
        "等价于：对任意 $c \\in \\mathbb{R}^p$",
        "$\\var(c^\\top \\tilde{\\beta}) \\geq \\var(c^\\top \\hat{\\beta})$",
        "每个坐标分量 $\\hat{\\beta}_j$ 方差最小",
        "对任意线性组合 $c^\\top \\tilde{\\beta}$ 方差不小于 $c^\\top \\hat{\\beta}$"
    ]
    create_slide_content(slide, "BLUE 的含义：协方差矩阵序", bullets, prs, "Theorem 4.2")

def create_why_linear_slide(prs):
    """Slide 12: Why only linear estimators?"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bullets = [
        "$A$ 可以是 $X$ 的任意非线性函数",
        "线性约束已包含极广的估计量类",
        "无偏性是自然要求",
        "在许多现代应用中，有偏估计",
        "(Ridge、Lasso) 方差更小"
    ]
    create_slide_content(slide, "为什么只比较线性估计量？", bullets, prs, "Intuition")

def create_projection_geometry_slide(prs):
    """Slide 13: OLS projection geometry."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bullets = [
        "投影矩阵 $H = X(X^\\top X)^{-1} X^\\top$",
        "$\\hat{Y} = HY$ 是 $Y$ 到 $X$ 列空间的投影",
        "残差 $\\hat{\\varepsilon} = (I - H)Y$ 与列空间正交",
        "$Y = \\hat{Y} + \\hat{\\varepsilon}$",
        "$H$ 和 $I_n - H$ 是投影矩阵，两者正交"
    ]
    create_slide_content(slide, "OLS 投影几何直观", bullets, prs, "Lemma 4.1")

def create_fitted_residual_slide(prs):
    """Slide 14: Fitted values and residuals distribution."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bullets = [
        "Theorem 4.2:",
        "$E(\\hat{Y}) = X\\beta$，$E(\\hat{\\varepsilon}) = 0$",
        "$\\cov(\\hat{Y}) = \\sigma^2 H$",
        "$\\cov(\\hat{\\varepsilon}) = \\sigma^2(I - H)$",
        "$\\hat{Y}$ 与 $\\hat{\\varepsilon}$ 不相关"
    ]
    create_slide_content(slide, "拟合值与残差的分布性质", bullets, prs, "Theorem 4.2")

def create_orthogonal_vs_uncorrelated_slide(prs):
    """Slide 15: Orthogonal vs uncorrelated."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bullets = [
        "陈述 (S1): $\\hat{Y}$ 与 $\\hat{\\varepsilon}$ 正交",
        "——代数事实，无需随机假设（OLS 投影性质）",
        "陈述 (S2): $\\hat{Y}$ 与 $\\hat{\\varepsilon}$ 不相关",
        "——随机陈述，需要 Gauss-Markov 假设",
        "两者含义不同，不能混淆"
    ]
    create_slide_content(slide, "正交 vs. 不相关", bullets, prs, "常见误区")

def create_error_ellipse_slide(prs):
    """Slide 16: Error ellipse visualization."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bullets = [
        "$\\hat{\\beta}$ 的置信椭圆",
        "其他线性无偏估计的置信椭圆",
        "OLS 的置信椭圆最小",
        "（被其他所有椭圆包裹）",
        "$\\Rightarrow$ $\\hat{\\beta}$ 在所有方向上方差最小"
    ]
    create_slide_content(slide, "误差椭圆可视化", bullets, prs, "Intuition")

def create_proof_outline_slide(prs):
    """Slide 17: Gauss-Markov proof outline."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bullets = [
        "第一步: OLS 满足无偏性条件 $AX = I_p$",
        "第二步: $\\cov(\\tilde{\\beta}) = \\cov(\\hat{\\beta}) + \\cov(\\tilde{\\beta} - \\hat{\\beta})$",
        "第三步: $\\cov(\\hat{\\beta}, \\tilde{\\beta} - \\hat{\\beta}) = 0$",
        "第四步: $\\cov(\\tilde{\\beta} - \\hat{\\beta}) \\succeq 0$",
        "$\\Rightarrow \\cov(\\tilde{\\beta}) \\succeq \\cov(\\hat{\\beta})$"
    ]
    create_slide_content(slide, "Gauss-Markov 定理证明思路", bullets, prs, "Proof")

def create_simple_regression_example_slide(prs):
    """Slide 18: Simple linear regression example."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bullets = [
        "$y_i = \\alpha + \\beta x_i + \\varepsilon_i$",
        "$\\var(\\hat{\\beta}) = \\sigma^2 / \\sum (x_i - \\bar{x})^2$",
        "任何线性无偏估计的 $\\beta$ 系数方差都不小于此值",
        "设计越分散，OLS 估计越精确"
    ]
    create_slide_content(slide, "例子：简单线性回归的 BLUE", bullets, prs, "Example")

def create_mean_blue_slide(prs):
    """Slide 19: Mean BLUE example."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bullets = [
        "$y_i \\sim$ 均值 $\\mu$，方差 $\\sigma^2$，互不相关",
        "线性估计 $\\hat{\\mu} = \\sum a_i y_i$，无偏要求 $\\sum a_i = 1$",
        "最优选择：$a_i = 1/n$（简单平均）",
        "方差 $\\var(\\hat{\\mu}) = \\sigma^2/n$ 为最小"
    ]
    create_slide_content(slide, "例子：均值的 BLUE", bullets, prs, "Example")

def create_prediction_theorem_slide(prs):
    """Slide 20: Gauss-Markov prediction theorem."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bullets = [
        "Gauss-Markov for Prediction:",
        "$\\hat{Y} = X\\hat{\\beta}$ 是 $X\\beta$ 的最佳线性无偏预测",
        "适用于任何线性预测 $\\tilde{Y} = \\tilde{H} Y$",
        "$\\cov(\\tilde{Y}) \\succeq \\cov(\\hat{Y})$"
    ]
    create_slide_content(slide, "例子：Gauss-Markov 预测定理", bullets, prs, "Theorem 4.3")

def create_summary_slide(prs):
    """Slide 21: Summary."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bullets = [
        "Gauss-Markov 模型: 线性 + 同方差 + 不相关",
        "OLS $\\hat{\\beta}$ 是 BLUE: $\\cov(\\tilde{\\beta}) \\succeq \\cov(\\hat{\\beta})$",
        "核心不等式: $\\cov(\\tilde{\\beta}) - \\cov(\\hat{\\beta}) = \\cov(\\tilde{\\beta} - \\hat{\\beta}) \\succeq 0$",
        "局限：",
        "  - 不谈非线性估计（Ridge、Lasso）",
        "  - 不谈正态假设（Normal 模型 MLE）",
        "  - 不谈稳健性"
    ]
    create_slide_content(slide, "Gauss-Markov 定理总结", bullets, prs, "Summary")

def main():
    # Create presentation with 16:9 aspect ratio
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # Generate slides
    create_cover_slide(prs)
    create_motivation_slide_1(prs)
    create_motivation_slide_2(prs)
    create_transition_slide(prs)
    create_assumptions_slide(prs)
    create_individual_level_slide(prs)
    create_homoskedasticity_slide(prs)
    create_ols_form_slide(prs)
    create_ols_mean_var_slide(prs)
    create_gauss_markov_theorem_slide(prs)
    create_blue_meaning_slide(prs)
    create_why_linear_slide(prs)
    create_projection_geometry_slide(prs)
    create_fitted_residual_slide(prs)
    create_orthogonal_vs_uncorrelated_slide(prs)
    create_error_ellipse_slide(prs)
    create_proof_outline_slide(prs)
    create_simple_regression_example_slide(prs)
    create_mean_blue_slide(prs)
    create_prediction_theorem_slide(prs)
    create_summary_slide(prs)

    # Save
    prs.save(OUTPUT_PATH)
    print(f"PPTX saved to: {OUTPUT_PATH}")

    # Verify file exists and size
    if os.path.exists(OUTPUT_PATH):
        size = os.path.getsize(OUTPUT_PATH)
        print(f"File size: {size} bytes ({size/1024:.1f} KB)")
        if size > 50000:
            print("✓ File size is reasonable (>50KB)")
        else:
            print("⚠ File size might be too small")
    else:
        print("✗ File was not created!")

if __name__ == "__main__":
    main()
